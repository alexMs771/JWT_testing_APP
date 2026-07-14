"""
Generate a PowerPoint deck for the cf-token-manager POC.
Run:  python build_poc_ppt.py
Output: POC-cf-token-manager.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------------------------------
# Palette (matches the app's hero gradient: #0d1b3e -> #17408b -> #2d6cdf)
# ----------------------------------------------------------------------------
NAVY      = RGBColor(0x0D, 0x1B, 0x3E)
BLUE      = RGBColor(0x17, 0x40, 0x8B)
BRIGHT    = RGBColor(0x2D, 0x6C, 0xDF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xF3, 0xF6, 0xFC)
GREY      = RGBColor(0x5B, 0x66, 0x77)
DARK      = RGBColor(0x1A, 0x22, 0x33)
GREEN     = RGBColor(0x19, 0x87, 0x54)
CODE_BG   = RGBColor(0x0E, 0x17, 0x2A)
CODE_FG   = RGBColor(0xE6, 0xED, 0xF7)

FONT      = "Segoe UI"
MONO      = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# helpers
# ----------------------------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def gradient_band(s, x, y, w, h):
    """Fake a diagonal gradient with three stacked rectangles."""
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.line.fill.background()
    shp.shadow.inherit = False
    fill = shp.fill
    fill.gradient()
    try:
        stops = fill.gradient_stops
        stops[0].color.rgb = NAVY
        stops[0].position = 0.0
        stops[1].color.rgb = BRIGHT
        stops[1].position = 1.0
        fill.gradient_angle = 35.0
    except Exception:
        fill.solid()
        fill.fore_color.rgb = BLUE
    return shp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, size, bold, color, font)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (text, size, bold, color, font) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = font
            r.font.color.rgb = color
    return tb


def bullet(text, size=15, bold=False, color=DARK, font=FONT):
    return [(text, size, bold, color, font)]


def header(s, kicker, title):
    """Standard content-slide header with a top accent bar."""
    rect(s, 0, 0, SW, Inches(0.28), BRIGHT)
    txt(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
        [[(kicker.upper(), 12, True, BRIGHT, FONT)]])
    txt(s, Inches(0.7), Inches(0.82), Inches(12), Inches(0.8),
        [[(title, 30, True, NAVY, FONT)]])
    rect(s, Inches(0.72), Inches(1.55), Inches(1.1), Pt(3), BRIGHT)


def footer(s, idx):
    txt(s, Inches(0.7), Inches(7.05), Inches(8), Inches(0.3),
        [[("cf-token-manager  ·  Proof of Concept  ·  v1.0.0", 9, False, GREY, FONT)]])
    txt(s, Inches(12.1), Inches(7.05), Inches(0.8), Inches(0.3),
        [[(str(idx), 9, True, GREY, FONT)]], align=PP_ALIGN.RIGHT)


def code_box(s, x, y, w, h, lines, size=12):
    rect(s, x, y, w, h, CODE_BG)
    tb = s.shapes.add_textbox(x + Inches(0.15), y + Inches(0.12),
                              w - Inches(0.3), h - Inches(0.24))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        p.line_spacing = 1.05
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.name = MONO
        r.font.color.rgb = CODE_FG
    return tb


def card(s, x, y, w, h, title, body_paras, title_color=NAVY, accent=BRIGHT):
    rect(s, x, y, w, h, WHITE, line=RGBColor(0xDD, 0xE3, 0xEC))
    rect(s, x, y, Inches(0.09), h, accent)
    tb = s.shapes.add_textbox(x + Inches(0.28), y + Inches(0.18),
                              w - Inches(0.5), h - Inches(0.32))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_after = Pt(6)
    r = p.add_run(); r.text = title
    r.font.size = Pt(16); r.font.bold = True; r.font.name = FONT
    r.font.color.rgb = title_color
    for para in body_paras:
        p = tf.add_paragraph()
        p.space_after = Pt(3)
        p.line_spacing = 1.05
        for (text, size, bold, color) in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.bold = bold
            r.font.name = FONT; r.font.color.rgb = color
    return tb


def table(s, x, y, w, rows, col_widths, header_fill=NAVY, fs=12, row_h=0.42):
    ncols = len(rows[0])
    nrows = len(rows)
    gt = s.shapes.add_table(nrows, ncols, x, y, w,
                            Inches(row_h * nrows)).table
    gt.first_row = True
    gt.horz_banding = True
    total = sum(col_widths)
    for ci, cw in enumerate(col_widths):
        gt.columns[ci].width = Emu(int(w * cw / total))
    for ri, row in enumerate(rows):
        gt.rows[ri].height = Inches(row_h)
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.size = Pt(fs)
            r.font.name = FONT
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
                r.font.color.rgb = WHITE; r.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                r.font.color.rgb = DARK
    return gt


def pill(s, x, y, w, text, color=BRIGHT, fg=WHITE, h=0.42, fs=12):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = True; r.font.name = FONT
    r.font.color.rgb = fg
    return shp


# ============================================================================
# SLIDE 1 — Title
# ============================================================================
s = slide()
gradient_band(s, 0, 0, SW, SH)
rect(s, 0, Inches(5.55), SW, Pt(4), BRIGHT)
txt(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(0.5),
    [[("PROOF OF CONCEPT", 16, True, RGBColor(0x9F, 0xC0, 0xF5), FONT)]])
txt(s, Inches(0.9), Inches(1.75), Inches(11.6), Inches(1.6),
    [[("cf-token-manager", 54, True, WHITE, FONT)]])
txt(s, Inches(0.9), Inches(3.05), Inches(11.4), Inches(1.0),
    [[("A JWT Access & Refresh Token Lifecycle Package", 24, False, WHITE, FONT)],
     [("for CFML / ColdBox", 24, False, RGBColor(0xC9, 0xDA, 0xF8), FONT)]],
    space_after=2)
# meta pills
pill(s, Inches(0.9), Inches(4.7), Inches(2.5), "CommandBox + ColdBox",
     color=RGBColor(0x1E, 0x4E, 0xA6), h=0.5, fs=13)
pill(s, Inches(3.6), Inches(4.7), Inches(2.2), "Adobe · Lucee · BoxLang",
     color=RGBColor(0x1E, 0x4E, 0xA6), h=0.5, fs=13)
pill(s, Inches(6.0), Inches(4.7), Inches(1.7), "Version 1.0.0",
     color=RGBColor(0x1E, 0x4E, 0xA6), h=0.5, fs=13)
pill(s, Inches(7.9), Inches(4.7), Inches(1.7), "MIT · ForgeBox",
     color=RGBColor(0x1E, 0x4E, 0xA6), h=0.5, fs=13)
txt(s, Inches(0.9), Inches(5.9), Inches(11), Inches(1.0),
    [[("ColdFusion Team  ·  MitrahSoft", 15, True, WHITE, FONT)],
     [("Validated against a live ColdBox host application", 13, False,
       RGBColor(0xC9, 0xDA, 0xF8), FONT)]], space_after=3)


# ============================================================================
# SLIDE 2 — Problem & Motivation
# ============================================================================
s = slide()
header(s, "Why we built it", "The Problem & Our Idea")
card(s, Inches(0.7), Inches(1.9), Inches(5.85), Inches(2.15),
     "The gap we found on ForgeBox",
     [[("Many JWT packages exist \u2014 but they stop at low-level", 14, False, DARK)],
      [("encode / decode / verify primitives.", 14, True, BLUE)],
      [("None package the access + refresh token concept as a", 14, False, DARK)],
      [("first-class, developer-ready API.", 14, True, BLUE)],
      [("Every app re-implements silent refresh & logout by hand.", 14, False, GREY)]],
     accent=RGBColor(0xE0, 0x7A, 0x3B))
card(s, Inches(6.75), Inches(1.9), Inches(5.85), Inches(2.15),
     "Our idea",
     [[("Reuse a proven low-level JWT engine for the crypto.", 14, False, DARK)],
      [("Build a lifecycle layer on top: access + refresh,", 14, False, DARK)],
      [("verify, decode & refresh \u2014 standardised & opinionated.", 14, False, DARK)],
      [("Sign via Java libs \u2192 runs on every CF engine.", 14, False, DARK)],
      [("Ship v1 now; role-based auth in v2.", 14, True, GREEN)]],
     accent=GREEN)
# task strip
rect(s, Inches(0.7), Inches(4.35), Inches(11.9), Inches(1.75), LIGHT)
rect(s, Inches(0.7), Inches(4.35), Inches(0.12), Inches(1.75), BRIGHT)
txt(s, Inches(1.0), Inches(4.55), Inches(11.2), Inches(0.4),
    [[("THE TASK", 12, True, BRIGHT, FONT)]])
txt(s, Inches(1.0), Inches(4.9), Inches(11.3), Inches(1.1),
    [[("Team lead assigned a POC: build a reusable JWT package via CommandBox/ColdBox and "
       "contribute it to ForgeBox. We analysed the existing packages, found the access/refresh "
       "gap, and designed cf-token-manager to fill it with a clean, standard API for developers.",
       15, False, DARK, FONT)]], line_spacing=1.15)
footer(s, 2)


# ============================================================================
# SLIDE 3 — Architecture
# ============================================================================
s = slide()
header(s, "How it's built", "Two-Layer Architecture")
# layer boxes
card(s, Inches(0.7), Inches(1.95), Inches(6.0), Inches(1.45),
     "Host ColdBox Application",
     [[("inject=\"TokenManager@cf-token-manager\"", 13, False, GREY)],
      [("The app only ever talks to the public API.", 13, False, DARK)]],
     accent=GREY)
card(s, Inches(0.7), Inches(3.55), Inches(6.0), Inches(1.6),
     "TokenManager.cfc  \u2014  public lifecycle API",
     [[("issue() · verify() · decode() · refresh()", 13, True, BLUE)],
      [("validate() · diagnose()", 13, True, BLUE)],
      [("Owns policy: TTLs, access vs refresh, iat/exp/iss.", 13, False, DARK)]],
     accent=BRIGHT)
card(s, Inches(0.7), Inches(5.3), Inches(6.0), Inches(1.5),
     "JWTService.cfc  \u2014  internal engine",
     [[("encode() · decode() · verifySignature()", 13, True, BLUE)],
      [("Java HMAC · base64url · constant-time compare.", 13, False, DARK)],
      [("Not injected by host apps.", 13, False, GREY)]],
     accent=NAVY)
# arrows
for yy in (3.35, 5.1):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(3.4), Inches(yy),
                           Inches(0.55), Inches(0.22))
    a.fill.solid(); a.fill.fore_color.rgb = BRIGHT
    a.line.fill.background(); a.shadow.inherit = False
# right column — why
card(s, Inches(6.95), Inches(1.95), Inches(5.65), Inches(4.85),
     "Why two layers?",
     [[("", 6, False, DARK)],
      [("\u2022  Engine can be upgraded without touching app code.", 15, False, DARK)],
      [("", 6, False, DARK)],
      [("\u2022  The lifecycle layer encodes all the \u201ccorrect\u201d", 15, False, DARK)],
      [("   behaviour so developers don't repeat it.", 15, False, DARK)],
      [("", 6, False, DARK)],
      [("\u2022  Reserved claims (type, iat, exp, iss) are owned", 15, False, DARK)],
      [("   by the package \u2014 callers can't forge type/expiry.", 15, False, DARK)],
      [("", 6, False, DARK)],
      [("\u2022  Signature is verified before any claim is trusted.", 15, False, DARK)],
      [("", 6, False, DARK)],
      [("\u2022  Java crypto \u2192 identical on Adobe CF, Lucee, BoxLang.", 15, True, BLUE)]],
     accent=GREEN)
footer(s, 3)


# ============================================================================
# SLIDE 4 — Public API
# ============================================================================
s = slide()
header(s, "The developer experience", "Public API \u2014 4 Core Functions")
apis = [
    ("issue()", "Mint an access + refresh token pair for a subject.", BRIGHT),
    ("verify()", "Boolean: authentic + unexpired + correct type. Never throws.", GREEN),
    ("decode()", "Read the claims struct (no signature check).", RGBColor(0xE0, 0x7A, 0x3B)),
    ("refresh()", "Swap a valid refresh token for a new access token.", NAVY),
]
x = Inches(0.7)
for name, desc, col in apis:
    card(s, x, Inches(1.9), Inches(2.85), Inches(1.85), name,
         [[(desc, 13, False, DARK)]], title_color=col, accent=col)
    x += Inches(3.0)
code_box(s, Inches(0.7), Inches(4.0), Inches(7.4), Inches(2.55),
    [ "// inject anywhere",
      'property name="tokenManager" inject="TokenManager@cf-token-manager";',
      "",
      "// login",
      'var tokens = tokenManager.issue({ id:101, role:"admin" });',
      "//  { accessToken:\"eyJ..\",  refreshToken:\"eyJ..\" }",
      "",
      "// protected route",
      "if ( tokenManager.verify( tokens.accessToken ) ) { ... }",
      "",
      "// access token expired",
      "var fresh = tokenManager.refresh( tokens.refreshToken );" ],
    size=12)
card(s, Inches(8.35), Inches(4.0), Inches(4.25), Inches(2.55),
     "Smart defaults",
     [[("id \u2192 becomes the JWT sub claim.", 13, False, DARK)],
      [("Other keys (role, email) carried verbatim.", 13, False, DARK)],
      [("Claim names lower-cased for cross-engine", 13, False, DARK)],
      [("consistency.", 13, False, DARK)],
      [("Missing id/sub \u2192 InvalidSubject thrown.", 13, True, BLUE)]],
     accent=BRIGHT)
footer(s, 4)


# ============================================================================
# SLIDE 5 — Diagnostics / error codes
# ============================================================================
s = slide()
header(s, "Precise error reporting", "Diagnostics \u2014 validate() & diagnose()")
txt(s, Inches(0.7), Inches(1.75), Inches(12), Inches(0.5),
    [[("verify() gives a boolean. When you need to tell the client WHY a token failed, "
       "diagnose() returns a code (non-throwing) and validate() throws a typed exception.",
       14, False, GREY, FONT)]], line_spacing=1.1)
rows = [
    ["diagnose() code", "Thrown exception", "Meaning"],
    ["VALID", "\u2014", "Authentic, unexpired, correct type"],
    ["MISSING_TOKEN", "cftokenmanager.MissingToken", "Empty / no token supplied"],
    ["MALFORMED_TOKEN", "cftokenmanager.MalformedToken", "Not a 3-segment JWT / unparseable body"],
    ["INVALID_SIGNATURE", "cftokenmanager.InvalidSignature", "Tampered or wrong secret"],
    ["TOKEN_EXPIRED", "cftokenmanager.ExpiredToken", "exp claim is in the past"],
    ["INVALID_TOKEN_TYPE", "cftokenmanager.InvalidTokenType", "Refresh token used on an access endpoint"],
]
table(s, Inches(0.7), Inches(2.55), Inches(11.9), rows,
      col_widths=[3, 4.5, 5], fs=13, row_h=0.55)
footer(s, 5)


# ============================================================================
# SLIDE 6 — Configuration & Security
# ============================================================================
s = slide()
header(s, "Policy & hardening", "Configuration & Security")
rows = [
    ["Setting", "Default", "Purpose"],
    ["secret", "override in prod", "HMAC signing key (from .env)"],
    ["algorithm", "HS512", "HS256 | HS384 | HS512"],
    ["accessTokenExpiry", "900 (15 min)", "Short-lived access token"],
    ["refreshTokenExpiry", "2592000 (30 d)", "Long-lived refresh token"],
    ["issuer", "cf-token-manager", "Stamped as the iss claim"],
]
table(s, Inches(0.7), Inches(2.0), Inches(6.2), rows,
      col_widths=[3.2, 3.0, 3.6], fs=12.5, row_h=0.55)
card(s, Inches(7.15), Inches(2.0), Inches(5.45), Inches(4.6),
     "Security design",
     [[("\u2022  Java javax.crypto HMAC \u2014 portable & proven.", 14, False, DARK)],
      [("", 5, False, DARK)],
      [("\u2022  Constant-time signature compare (anti-timing).", 14, False, DARK)],
      [("", 5, False, DARK)],
      [("\u2022  base64url per JWT spec (URL-safe, no pad).", 14, False, DARK)],
      [("", 5, False, DARK)],
      [("\u2022  Reserved claims stripped from caller input.", 14, False, DARK)],
      [("", 5, False, DARK)],
      [("\u2022  Signature verified BEFORE claims are read.", 14, False, DARK)],
      [("", 5, False, DARK)],
      [("\u2022  Epoch cast to Java long \u2192 clean NumericDate.", 14, False, DARK)]],
     accent=GREEN)
footer(s, 6)


# ============================================================================
# SLIDE 7 — How it was tested
# ============================================================================
s = slide()
header(s, "Validation", "How the POC Was Tested")
card(s, Inches(0.7), Inches(1.95), Inches(5.85), Inches(2.5),
     "1 · Engineering self-test  \u2192  /demo",
     [[("TokenDemo.cfc renders a self-test matrix on every load.", 13, False, DARK)],
      [("Runs expected-vs-actual for every function across all", 13, False, DARK)],
      [("success AND failure dimensions: tampered, wrong secret,", 13, False, DARK)],
      [("expired, malformed, missing, wrong type.", 13, False, DARK)],
      [("Interactive issue / refresh / inspect console.", 13, True, BLUE)]],
     accent=BRIGHT)
card(s, Inches(6.75), Inches(1.95), Inches(5.85), Inches(2.5),
     "2 · Real auth flow  \u2192  login/dashboard",
     [[("Security.cfc uses the package exactly like production:", 13, False, DARK)],
      [("Login \u2192 issue() mints the first pair.", 13, False, DARK)],
      [("Dashboard \u2192 decode() drives live countdowns.", 13, False, DARK)],
      [("Access expires \u2192 refresh() silently renews.", 13, False, DARK)],
      [("Refresh expires \u2192 auto logout.", 13, True, BLUE)]],
     accent=GREEN)
rect(s, Inches(0.7), Inches(4.75), Inches(11.9), Inches(1.85), LIGHT)
rect(s, Inches(0.7), Inches(4.75), Inches(0.12), Inches(1.85), RGBColor(0xE0, 0x7A, 0x3B))
txt(s, Inches(1.0), Inches(4.95), Inches(11.3), Inches(0.4),
    [[("DEMO TTLs (so the whole lifecycle is visible in ~2 minutes)", 12, True,
       RGBColor(0xC0, 0x60, 0x20), FONT)]])
code_box(s, Inches(1.0), Inches(5.4), Inches(6.6), Inches(1.05),
    [ '"cf-token-manager" : {',
      '    accessTokenExpiry  : 30,   // 30 seconds',
      '    refreshTokenExpiry : 120   // 2 minutes',
      "}" ], size=12)
txt(s, Inches(7.9), Inches(5.5), Inches(4.5), Inches(1.0),
    [[("Production values", 13, True, NAVY, FONT)],
     [("900 / 2,592,000 seconds", 13, False, DARK, FONT)],
     [("(15 min access · 30 day refresh)", 12, False, GREY, FONT)]], space_after=3)
footer(s, 7)


# ============================================================================
# SLIDE 8 — Results
# ============================================================================
s = slide()
header(s, "Outcome", "Test Results \u2014 All Checks Pass")
rows = [
    ["Scenario", "Function", "Expected", "Result"],
    ["Mint valid pair", "issue()", "two 3-seg JWTs", "PASS"],
    ["Reject subject w/o id", "issue()", "InvalidSubject", "PASS"],
    ["Read claims unverified", "decode()", "claims returned", "PASS"],
    ["Valid access token", "verify()", "true", "PASS"],
    ["Tampered / wrong secret", "diagnose()", "INVALID_SIGNATURE", "PASS"],
    ["Expired token", "diagnose()", "TOKEN_EXPIRED", "PASS"],
    ["Malformed / missing", "diagnose()", "MALFORMED / MISSING", "PASS"],
    ["Wrong type (access as refresh)", "diagnose()", "INVALID_TOKEN_TYPE", "PASS"],
    ["Valid refresh \u2192 new access", "refresh()", "new token VALID", "PASS"],
    ["Access used as refresh", "refresh()", "InvalidTokenType", "PASS"],
]
table(s, Inches(0.7), Inches(1.95), Inches(9.1), rows,
      col_widths=[4.2, 2.2, 3.6, 1.6], fs=11.5, row_h=0.42)
# big pass badge
badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(10.1), Inches(2.6), Inches(2.5), Inches(2.4))
badge.fill.solid(); badge.fill.fore_color.rgb = GREEN
badge.line.fill.background(); badge.shadow.inherit = False
tf = badge.text_frame; tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "12 / 12"
r.font.size = Pt(34); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r = p2.add_run(); r.text = "checks passed"
r.font.size = Pt(15); r.font.name = FONT; r.font.color.rgb = WHITE
txt(s, Inches(10.1), Inches(5.15), Inches(2.5), Inches(1.0),
    [[("Matrix re-runs live on", 11, False, GREY, FONT)],
     [("every page load \u2014", 11, False, GREY, FONT)],
     [("regressions show instantly.", 11, False, GREY, FONT)]],
    align=PP_ALIGN.CENTER, space_after=1)
footer(s, 8)


# ============================================================================
# SLIDE 9 — Roadmap
# ============================================================================
s = slide()
header(s, "What's next", "Roadmap")
# v1
card(s, Inches(0.7), Inches(2.0), Inches(3.75), Inches(4.4),
     "v1.0.0  \u2014  this POC",
     [[("Access + refresh lifecycle", 13, True, GREEN)],
      [("issue / verify / decode / refresh", 13, False, DARK)],
      [("validate / diagnose", 13, False, DARK)],
      [("Configurable policy", 13, False, DARK)],
      [("HS256 / 384 / 512", 13, False, DARK)],
      [("Adobe · Lucee · BoxLang", 13, False, DARK)],
      [("", 8, False, DARK)],
      [("STATUS: validated \u2713", 13, True, GREEN)]],
     accent=GREEN)
# v2
card(s, Inches(4.75), Inches(2.0), Inches(3.75), Inches(4.4),
     "v2  \u2014  planned",
     [[("Role-based authentication", 13, True, BRIGHT)],
      [("Roles: user, admin,", 13, False, DARK)],
      [("subscribed users", 13, False, DARK)],
      [("Guard routes / actions by role", 13, False, DARK)],
      [("More standards & helpers", 13, False, DARK)],
      [("", 8, False, DARK)],
      [("STATUS: next up", 13, True, BRIGHT)]],
     accent=BRIGHT)
# future
card(s, Inches(8.8), Inches(2.0), Inches(3.8), Inches(4.4),
     "Future",
     [[("Token revocation / blocklist", 13, False, DARK)],
      [("Refresh-token rotation store", 13, False, DARK)],
      [("Asymmetric algorithms", 13, False, DARK)],
      [("(RS / ES families)", 13, False, DARK)],
      [("", 8, False, DARK)],
      [("ForgeBox public release", 13, True, NAVY)]],
     accent=NAVY)
footer(s, 9)


# ============================================================================
# SLIDE 10 — Conclusion
# ============================================================================
s = slide()
gradient_band(s, 0, 0, SW, SH)
rect(s, 0, Inches(0), SW, Pt(6), BRIGHT)
txt(s, Inches(0.9), Inches(0.9), Inches(11.5), Inches(0.5),
    [[("CONCLUSION", 15, True, RGBColor(0x9F, 0xC0, 0xF5), FONT)]])
txt(s, Inches(0.9), Inches(1.45), Inches(11.6), Inches(1.6),
    [[("A real gap, cleanly closed.", 38, True, WHITE, FONT)]])
txt(s, Inches(0.9), Inches(2.7), Inches(11.4), Inches(1.6),
    [[("cf-token-manager delivers a standardised access & refresh token lifecycle for "
       "CFML/ColdBox \u2014 built on a proven JWT engine, portable across every CF engine via "
       "Java crypto, and validated end-to-end.", 18, False,
       RGBColor(0xE7, 0xEF, 0xFB), FONT)]], line_spacing=1.2)
# recommendation box
rect(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(1.5),
     RGBColor(0x12, 0x2C, 0x5E))
rect(s, Inches(0.9), Inches(4.5), Inches(0.14), Inches(1.5), GREEN)
txt(s, Inches(1.25), Inches(4.72), Inches(11), Inches(0.4),
    [[("RECOMMENDATION", 13, True, RGBColor(0x8F, 0xD9, 0xAE), FONT)]])
txt(s, Inches(1.25), Inches(5.12), Inches(11), Inches(0.8),
    [[("Polish and publish Version 1 to ForgeBox, then begin Version 2 "
       "(role-based authentication).", 17, True, WHITE, FONT)]], line_spacing=1.1)
txt(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
    [[("Tested with the TestingApp ColdBox application  ·  Adobe ColdFusion 2025 via CommandBox",
       12, False, RGBColor(0xC9, 0xDA, 0xF8), FONT)]])


prs.save("POC-cf-token-manager.pptx")
print("Saved POC-cf-token-manager.pptx with", len(prs.slides._sldIdLst), "slides")
